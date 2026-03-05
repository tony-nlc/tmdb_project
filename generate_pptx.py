# ============================================
# PowerPoint Generator for TMDB Deliverable 3
# Run this in your local Python environment
# ============================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
GOLD = RGBColor(218, 165, 32)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.8), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    return slide

# ============================================
# SLIDE 1: Title
# ============================================
add_title_slide(prs, 
    "TMDB Movie Data Analysis",
    "Deliverable 3: EDA, Statistical Inference & Tableau Dashboard")

# ============================================
# SLIDE 2: Overview
# ============================================
add_content_slide(prs, "Project Overview", [
    "Dataset: TMDB (The Movie Database) - 3,901 movies with budget & revenue",
    "Features: Budget, Revenue, Ratings, Genres, Release Dates, Runtime",
    "Tools Used: Python (Pandas, Matplotlib, Seaborn, SciPy), Tableau",
    "Objective: Derive actionable insights from movie data through EDA & statistics"
])

# ============================================
# SLIDE 3: EDA
# ============================================
add_content_slide(prs, "Exploratory Data Analysis", [
    "Total Records: 3,901 movies (with budget & revenue > 0)",
    "Features Analyzed: Budget, Revenue, Runtime, Popularity, Ratings",
    "Data Quality: Handled missing values, created log transformations",
    "Top Genres by Volume: Drama (1,812), Comedy (1,479), Action (1,003)",
    "Top by Revenue: Animation ($262M avg), Adventure ($227M), Action ($173M)"
])

# ============================================
# SLIDE 4: Visualizations
# ============================================
add_content_slide(prs, "EDA Visualizations Created", [
    "Histograms: Distribution of budget, revenue, runtime, popularity",
    "Box Plots: Identify outliers in key metrics by genre",
    "Bar Charts: Top genres, top movie franchises by revenue",
    "Scatter Plots: Budget vs Revenue with trendlines",
    "Heatmaps: Correlation between numerical features"
])

# ============================================
# SLIDE 5: Confidence Intervals
# ============================================
add_content_slide(prs, "Statistical Inference: Confidence Intervals (95%)", [
    "Mean Revenue: $108,143,432",
    "   95% CI: [$102.6M, $113.7M]",
    "",
    "Mean Budget: $36,709,655",
    "   95% CI: [$35.4M, $38.0M]",
    "",
    "Mean Vote Average: 6.46",
    "   95% CI: [6.43, 6.48]"
])

# ============================================
# SLIDE 6: Hypothesis Tests
# ============================================
add_two_column_slide(prs, "Hypothesis Testing Results",
    "H1: Action vs Drama Budgets",
    [
        "H0: Mean budget(Action) = Mean budget(Drama)",
        "H1: Mean budget(Action) > Mean budget(Drama)",
        "Result: ✅ Reject H0 (p < 0.001)",
        "Action (n=1,003): Avg = $60.9M",
        "Drama (n=1,812): Avg = $26.8M",
        "t-statistic: 21.41"
    ],
    "H2: Homepage Impact",
    [
        "H0: Revenue(with) = Revenue(without)",
        "H1: Revenue(with) > Revenue(without)",
        "Result: ✅ Reject H0 (p < 0.001)",
        "With homepage: $161.9M avg",
        "Without: $69.5M avg",
        "t-statistic: 16.63"
    ])

# ============================================
# SLIDE 7: Correlation Test
# ============================================
add_content_slide(prs, "Hypothesis Test: Budget-Revenue Correlation", [
    "H0: No correlation between budget and revenue (ρ = 0)",
    "H1: Significant correlation exists (ρ ≠ 0)",
    "",
    "Results:",
    "• Pearson Correlation: r = 0.71 (strong positive)",
    "• p-value: < 0.001 (highly significant)",
    "• Conclusion: Budget strongly predicts revenue",
    "• Interpretation: For every $1 increase in budget, revenue increases significantly"
])

# ============================================
# SLIDE 8: Tableau Dashboard
# ============================================
add_content_slide(prs, "Tableau Dashboard Features", [
    "Data Files: 6 CSV files generated (tableau_*.csv)",
    "• tableau_genre_summary.csv - Genre breakdown with avg revenue",
    "• tableau_year_trends.csv - Year-by-year trends (1916-2020)",
    "• tableau_budget_revenue.csv - For scatter plots (3,901 movies)",
    "• tableau_month_genre.csv - Revenue by month × genre",
    "• tableau_franchises.csv - Top 20 franchises",
    "• tableau_full_data.csv - Full cleaned dataset",
    "Interactive: Year slider, Genre dropdown, Budget range filter"
])

# ============================================
# SLIDE 9: Key Insights
# ============================================
add_content_slide(prs, "Key Insights & Recommendations", [
    "1. Budget matters: r = 0.71 correlation with revenue",
    "2. Top Genres by Revenue: Animation ($262M), Adventure ($227M), Family ($187M)",
    "3. Top Franchises: Avatar ($2.9B), Avengers ($1.5B), Frozen ($1.3B)",
    "4. Homepage Impact: Movies with homepages earn 2.3x more revenue",
    "5. Action vs Drama: Action movies have 2.3x higher budgets"
])

# ============================================
# SLIDE 10: Conclusion
# ============================================
add_content_slide(prs, "Conclusion", [
    "✅ Comprehensive EDA with multiple visualization types",
    "✅ Statistical inference with 95% confidence intervals",
    "✅ 3 hypothesis tests all statistically significant (p < 0.001)",
    "✅ Interactive Tableau dashboard with 6+ charts & filters",
    "Key Takeaway: Budget is the strongest predictor of movie success"
])

# ============================================
# SLIDE 11: Thank You
# ============================================
add_title_slide(prs, "Thank You!", "Questions?")

# Save
prs.save("TMDB_Deliverable3_Presentation.pptx")
print("✅ PowerPoint saved: TMDB_Deliverable3_Presentation.pptx")
